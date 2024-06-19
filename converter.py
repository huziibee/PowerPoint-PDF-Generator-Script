import os
import comtypes.client
from pptx import Presentation
import pandas as pd
from colorama import init, Fore, Style

init()




def replace_text_in_slide(slide, old_text, new_text):
    # Replaces old_text with new_text in a slide
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            if old_text in shape.text:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        if old_text in run.text:
                            run.text = run.text.replace(old_text, new_text)
                            
                            

def create_pdf_from_ppt(ppt_filename, pdf_filename):
    # Convert a PowerPoint file to PDF
    
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1
    
    # print(f"Creating {pdf_filename} from {ppt_filename}...")
    

    deck = powerpoint.Presentations.Open(ppt_filename)
    deck.SaveAs(pdf_filename, 32)  # 32 for PDF format
    deck.Close()
    powerpoint.Quit()
    name = pdf_filename.split("\\")[-1]
    print("New file created: "+ Style.BRIGHT + Fore.CYAN +name + Style.RESET_ALL)
    
    
def capitalize_name(full_name):
    return ' '.join(word.capitalize() for word in full_name.split())
    
def process_names(names_list, pptx_file_path):
    # Process a list of names, creating a PDF for each using the appropriate template based on name length
    
    short_template_path = "short name.pptx"  
    long_template_path = "long name.pptx"  
    output_folder = "output_pdfs" 
    name_length_threshold=11
    
    # Creates the output folder if it doesn't exist
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    for name in names_list:
        
        # Chooses the appropriate template based on the length of the name
        if len(name) >= name_length_threshold:
            if len(name) < 13 and ('I' in name.upper() or 'l' not in name):
                template_path = short_template_path
            else:
                template_path = long_template_path
        else:
            template_path = short_template_path
            
        template_path = pptx_file_path

        # Load the PowerPoint template
        prs = Presentation(template_path)
        
        # parts = name.split()

        # # Check if the name has more than one space (more than two parts)
        # if len(parts) > 2:
        #     print(f"Can't create for {name}: more than one space in the name.")
        #     continue 

        slide = prs.slides[0]
        replace_text_in_slide(slide, "[NAME]", capitalize_name(name))

        abs_output_folder = os.path.abspath(output_folder)
        temp_ppt_path = os.path.join(abs_output_folder, f"{name}.pptx")
        pdf_path = os.path.join(abs_output_folder, f"{name}.pdf")
        prs.save(temp_ppt_path)

        # Convert to PDF
        create_pdf_from_ppt(temp_ppt_path, pdf_path)

        # Remove the temporary PowerPoint file
        os.remove(temp_ppt_path)
        
        
### RUNNER CODE ###

# print("Welcome to the Name to PDF Converter!" + Style.RESET_ALL)

# check that we are in the ppt test aws folder 

if not os.getcwd().endswith("ppt test aws"):
    print(Fore.RED + "Please run this script in the ppt test aws folder" + Style.RESET_ALL)
    exit()

        
# Get the current working directory
current_directory = os.getcwd()
        
# List all files in the current directory
all_files = os.listdir(current_directory)

# Filter and list only the CSV files
csv_files = [file for file in all_files if file.endswith('.csv')]

pptx_files = [file for file in all_files if file.endswith('.pptx')]

print(Style.BRIGHT + Fore.BLUE+ "\nCSV files found in current directory:" + Style.RESET_ALL)

# Print the list of CSV files
for i in range(len(csv_files)):
    print(Style.BRIGHT+  f"{i}. {csv_files[i]}")
    

index1 = int(input(Style.BRIGHT + Fore.BLUE+"\nPlease enter the index for the CSV file:\n"+ Style.RESET_ALL + "eg. 0 (to choose Cloud_101_with_AWS.csv)\n-->"))
csv_file_path = os.path.join(os.getcwd(), csv_files[index1])

# Check if the file exists
if not os.path.exists(csv_file_path):
    print(f"File not found: {csv_file_path}")
    exit()



print(Style.BRIGHT + Fore.BLUE+"\nPPTX files found in current directory:"+ Style.RESET_ALL)

# Print the list of CSV files
for i in range(len(pptx_files)):
    print(Style.BRIGHT+ f"{i}. {pptx_files[i]}"+ Style.RESET_ALL)
    

index2 = int(input(Style.BRIGHT + Fore.BLUE+"\nPlease enter the index for the PPTX file:\n"+ Style.RESET_ALL + "eg. 0 (to choose name.pptx)\n-->"))
pptx_file_path = pptx_files[index2]

# print (csv_file_path)

# print(pptx_file_path)

# Reads the CSV file and converts the names column to a list
data = pd.read_csv(csv_files[index1])
temp_names = data['Name'].tolist()
names = [name.strip() for name in temp_names]

# print(names)

print(Style.BRIGHT + Fore.GREEN+"Processing names..."+ Style.RESET_ALL)


process_names(names, pptx_file_path)

